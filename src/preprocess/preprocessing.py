import re, os
from PyPDF2 import PdfReader
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
import pandas as pd
import extract_msg
from typing import List
from src.config import WORD, PPT, EMAIL, EXCEL, PDF, MIN_WORDS, logger
from langchain_community.document_loaders import PyPDFLoader as pdf


def load_pdf(file_path: str) -> List:
    """
    Extracts text and generate chunks from the pdf document

    :param file_path: path of the pdf file
    :return: identified chunks from the pdf file
    """
    texts = []
    loader = pdf(file_path)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200,
                                                   separators=["\n\n", "\n", " ", ".", ","])
    pages = loader.load_and_split(text_splitter)
    for p in pages:
        texts.append(p.page_content)
    return texts


def load_doc(file_path: str) -> List:
    """
    Extracts text and generate chunks from the word document

    :param file_path: path of the word file
    :return: identified chunks from the word file
    """
    texts = []
    doc = Document(file_path)
    for para in doc.paragraphs:
        texts.append(para.text)
    return texts


def load_ppt(file_path: str, min_words: int = MIN_WORDS) -> List:
    """
    Extracts text and generate chunks from the ppt document based on the slides.

    :param file_path: path of the ppt file
    :param min_words: minimum number of words to consider as chunk.
    :return: identified chunks from the ppt slides.
    """
    prs = Presentation(file_path)
    texts = []
    slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text += " " + shape.text

        if len(slide_text.split()) > min_words:
            texts.append(slide_text)
            slide_text = ""
    return texts


def load_excel(file_path: str, file_type: str) -> str:
    """
    Extracts texts from the excel file.

    :param file_path: path of the excel file
    :param file_type: extension of the excel file
    :return: text from the excel file
    """

    if file_type in ['xlsx', 'xls']:
        df = pd.read_excel(file_path)
    else:
        df = pd.read_csv(file_path)

    texts = df.to_string()
    return texts


def load_email(file_path: str) -> str:
    """
    Extracts text from the email file.

    :param file_path: path of the email file
    :return: joined text from the subject and the body of the email.
    """
    msg = extract_msg.Message(file_path)
    subject = msg.subject
    body = msg.body

    texts = subject + body
    return texts


def load_files(data_dir: str) -> List:
    """
    Extracts the text from the files identified in the data_dir according to the document type.
    Returns a list of the identified chunks.

    :param data_dir: directory of the data files for the chunks extraction.
    :return: list of the identified chunks.
    """
    files = []
    counts = {"ppt": 0, "word": 0, "pdf": 0, "email": 0, "excel": 0, "out of scope": 0}

    for filename in os.listdir(data_dir):
        file_path = os.path.join(data_dir, filename)
        file_name = os.path.basename(filename)
        file_type = file_name.split('.')[-1].lower()

        texts = ""
        if file_type == PDF:
            texts = [[txt, file_name] for txt in load_pdf(file_path)]
            counts["pdf"] += 1
        elif file_type in WORD:
            texts = [[txt, file_name] for txt in load_doc(file_path)]
            counts["word"] += 1
        elif file_type == PPT:
            texts = [[txt, file_name] for txt in load_ppt(file_path)]
            counts["ppt"] += 1
        elif file_type == EMAIL:
            texts = load_email(file_path)
            counts["email"] += 1
        elif file_type in EXCEL:
            texts = load_excel(file_path, file_type)
            counts["excel"] += 1
        else:
            texts = None
            logger.info("Out of scope:", file_name)
            counts["out of scope"] += 1

        if texts:
            if isinstance(texts, list):
                files.extend(texts)
            else:
                files.append(texts)

    logger.info("Number of files processed: {}".format(counts))
    logger.info("Total number of passages extracted: {}".format(len(files)))
    return files
